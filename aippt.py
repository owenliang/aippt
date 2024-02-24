import os
import json 
import broadscope_bailian
from pptx import Presentation

# 调用阿里云百炼大模型API
def bailian_llm(query,history=[],user_stop_words=[]): 
    # 百炼鉴权相关，从环境变量读取
    access_key_id=os.environ.get("ACCESS_KEY_ID")
    access_key_secret=os.environ.get("ACCESS_KEY_SECRET")
    agent_key=os.environ.get("AGENT_KEY")
    app_id=os.environ.get("APP_ID")
    
    # 构造LLM输入
    messages=[{'role':'system','content':'You are a helpful assistant.'}]
    for hist in history:
        messages.append({'role':'user','content':hist[0]})
        messages.append({'role':'assistant','content':hist[1]})
    messages.append({'role':'user','content':query})
    
    # 调用API
    client=broadscope_bailian.AccessTokenClient(access_key_id=access_key_id, access_key_secret=access_key_secret,agent_key=agent_key)
    resp=broadscope_bailian.Completions(token=client.get_token()).create(app_id=app_id,messages=messages,result_format="message",stop=user_stop_words)
    return resp['Data']['Choices'][0]['Message']['Content']

# 生成PPT内容
def generate_ppt_content(topic,pages):
    # 输出格式
    output_format=json.dumps({
        "title":"example title",
        "pages":[
            {
                "title": "title for page 1",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                ],
            },
            {
                "title": "title for page 2",
                "content": [
                    {
                        "title": "title for paragraph 1",
                        "description": "detail for paragraph 1",
                    },
                    {
                        "title": "title for paragraph 2",
                        "description": "detail for paragraph 2",
                    },
                    {
                        "title": "title for paragraph 3",
                        "description": "detail for paragraph 3",
                    },
                ],
            },
        ],
    },ensure_ascii=True)
    
    # prompt
    prompt=f'''我要准备1个关于{topic}的PPT，要求一共写{pages}页，请你根据主题生成详细内容，不要省略。
    按这个JSON格式输出{output_format}，只能返回JSON，且JSON不要用```包裹，内容要用中文。'''

    #print(prompt)
    
    # 调用llm生成PPT内容
    ppt_content=json.loads(bailian_llm(prompt))
    return ppt_content

# 生成PPT文件
def generate_ppt_file(topic,ppt_content):
    ppt=Presentation()
    
    # PPT首页
    slide=ppt.slides.add_slide(ppt.slide_layouts[0]) # title&subtitle layout
    slide.placeholders[0].text=ppt_content['title']
    slide.placeholders[1].text="通义千问72B"
    
    # 内容页
    print('总共%d页...'%len(ppt_content['pages']))
    for i,page in enumerate(ppt_content['pages']):
        print('生成第%d页:%s'%(i+1,page['title']))
        slide=ppt.slides.add_slide(ppt.slide_layouts[1]) # title&content layout
        # 标题
        slide.placeholders[0].text=page['title']
        # 正文
        for sub_content in page['content']:
            print(sub_content)
            # 一级正文
            sub_title=slide.placeholders[1].text_frame.add_paragraph()
            sub_title.text,sub_title.level=sub_content['title'],1
            # 二级正文
            sub_description=slide.placeholders[1].text_frame.add_paragraph()
            sub_description.text,sub_description.level=sub_content['description'],2
    
    ppt.save('%s.pptx'%topic)

if __name__=='__main__':
    while True:
        # 输入需求
        topic=input('输入主题:')
        pages=int(input('输入页数:'))
        # 生成PPT内容
        ppt_content=generate_ppt_content(topic,pages)
        # 生成PPT文件
        generate_ppt_file(topic,ppt_content)